import os
import re
import numpy as np
import matplotlib.pyplot as plt
from pptx import Presentation

# === CONFIG ===
pptx_path = "LOW-POWER-FIGURES-v5.0.pptx"
output_dir = "confusion_matrices"
selected_slides = list(range(6, 11)) + list(range(12, 15))
os.makedirs(output_dir, exist_ok=True)

# === LOAD PRESENTATION ===
prs = Presentation(pptx_path)

def extract_table_data(table):
    """
    Extracts numeric values and potential headers from a PowerPoint table.
    Returns: (matrix, row_labels, col_labels)
    """
    # Convert table to list of lists
    raw = [[cell.text.strip() for cell in row.cells] for row in table.rows]

    # Determine header rows/cols
    first_row = raw[0]
    first_col = [r[0] for r in raw]

    # Check if first row/col are mostly text
    row_has_text = any(re.search(r"[A-Za-z]", c) for c in first_row)
    col_has_text = any(re.search(r"[A-Za-z]", c) for c in first_col)

    # Extract numeric matrix
    start_row = 1 if col_has_text else 0
    start_col = 1 if row_has_text else 0

    numeric_matrix = []
    for row in raw[start_row:]:
        values = []
        for cell in row[start_col:]:
            match = re.findall(r"[-+]?\d*\.\d+|\d+", cell)
            if match:
                values.append(float(match[0]))
            else:
                values.append(np.nan)
        numeric_matrix.append(values)

    # Extract labels
    row_labels = [r[0] for r in raw[start_row:]] if col_has_text else [f"R{i}" for i in range(len(numeric_matrix))]
    col_labels = raw[0][start_col:] if row_has_text else [f"C{i}" for i in range(len(numeric_matrix[0]))]

    return numeric_matrix, row_labels, col_labels

def plot_confusion_matrix(matrix, row_labels, col_labels, title, filename):
    """Plots and saves confusion matrix with headers."""
    matrix = np.array(matrix, dtype=float)
    fig, ax = plt.subplots(figsize=(4, 4))
    cax = ax.matshow(matrix, cmap="Blues")
    plt.title(title, pad=15)
    plt.colorbar(cax)

    # Axis setup
    ax.set_xticks(np.arange(len(col_labels)))
    ax.set_yticks(np.arange(len(row_labels)))
    ax.set_xticklabels(col_labels, rotation=45, ha="right")
    ax.set_yticklabels(row_labels)
    plt.xlabel("Predicted")
    plt.ylabel("Actual")

    # Annotate cells
    for (i, j), val in np.ndenumerate(matrix):
        if not np.isnan(val):
            ax.text(j, i, f"{val:.0f}", ha="center", va="center", color="black")

    plt.tight_layout()
    plt.savefig(filename, dpi=400)
    plt.close()
    print(f"✅ Saved labeled confusion matrix: {filename}")

# === PROCESS SELECTED SLIDES ===
for idx in selected_slides:
    try:
        slide = prs.slides[idx - 1]
        found = False
        for shape in slide.shapes:
            if shape.has_table:
                matrix, row_labels, col_labels = extract_table_data(shape.table)
                if 2 <= len(matrix) <= 10 and 2 <= len(matrix[0]) <= 10:
                    found = True
                    fname = os.path.join(output_dir, f"confusion_matrix_slide_{idx}.png")
                    plot_confusion_matrix(matrix, row_labels, col_labels, f"Confusion Matrix (Slide {idx})", fname)
        if not found:
            print(f"ℹ️ No numeric confusion-matrix-like tables on slide {idx}")
    except Exception as e:
        print(f"❌ Error on slide {idx}: {e}")

print("\n🎉 All detected and labeled confusion matrices saved as PNGs.")
